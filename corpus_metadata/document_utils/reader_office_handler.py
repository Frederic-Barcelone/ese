#!/usr/bin/env python3
"""
================================
Office Document Handler Module
================================
Script: corpus_metadata/document_utils/reader_office_handler.py

Purpose:
--------
Handle Microsoft Office documents (Word, Excel, PowerPoint) with two extraction modes:
- intro: First portion of content (10k chars for Word, 5 sheets for Excel, 10 slides for PPT)
- full: Complete content extraction

Supported formats:
------------------
- Word: .docx, .doc
- Excel: .xlsx, .xls
- PowerPoint: .pptx, .ppt

Dependencies:
-------------
- python-docx: For Word documents
- openpyxl: For Excel files
- python-pptx: For PowerPoint presentations
- python-doc2txt: For legacy .doc files (optional)

Author: Corpus Metadata System
Date: 2024
"""

import logging
from pathlib import Path
from typing import Dict, Any, Optional, List, Tuple
import zipfile
import xml.etree.ElementTree as ET

logger = logging.getLogger(__name__)

# Try to import Office handling libraries
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available. Install with: pip install python-docx")

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    logger.warning("openpyxl not available. Install with: pip install openpyxl")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. Install with: pip install python-pptx")

try:
    import doc2txt
    DOC2TXT_AVAILABLE = True
except ImportError:
    DOC2TXT_AVAILABLE = False
    logger.debug("doc2txt not available for legacy .doc files")


class BaseOfficeHandler:
    """Base class for Office document handlers"""
    
    def __init__(self):
        """Initialize handler"""
        self.intro_chars = 10000  # Characters for intro mode
    
    def _extract_basic_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract basic metadata from Office files using ZIP structure"""
        metadata = {}
        
        try:
            if zipfile.is_zipfile(file_path):
                with zipfile.ZipFile(file_path, 'r') as zf:
                    # Try to extract core properties
                    if 'docProps/core.xml' in zf.namelist():
                        try:
                            core_xml = zf.read('docProps/core.xml')
                            root = ET.fromstring(core_xml)
                            
                            # Common properties
                            ns = {
                                'cp': 'http://schemas.openxmlformats.org/package/2006/metadata/core-properties',
                                'dc': 'http://purl.org/dc/elements/1.1/',
                                'dcterms': 'http://purl.org/dc/terms/'
                            }
                            
                            # Extract creator
                            creator = root.find('.//dc:creator', ns)
                            if creator is not None and creator.text:
                                metadata['author'] = creator.text
                            
                            # Extract title
                            title = root.find('.//dc:title', ns)
                            if title is not None and title.text:
                                metadata['title'] = title.text
                            
                            # Extract subject
                            subject = root.find('.//dc:subject', ns)
                            if subject is not None and subject.text:
                                metadata['subject'] = subject.text
                            
                            # Extract created date
                            created = root.find('.//dcterms:created', ns)
                            if created is not None and created.text:
                                metadata['created'] = created.text
                            
                            # Extract modified date
                            modified = root.find('.//dcterms:modified', ns)
                            if modified is not None and modified.text:
                                metadata['modified'] = modified.text
                                
                        except Exception as e:
                            logger.debug(f"Could not parse core properties: {e}")
                    
                    # Try to extract app properties
                    if 'docProps/app.xml' in zf.namelist():
                        try:
                            app_xml = zf.read('docProps/app.xml')
                            root = ET.fromstring(app_xml)
                            
                            # Extract application
                            app = root.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/extended-properties}Application')
                            if app is not None and app.text:
                                metadata['application'] = app.text
                                
                        except Exception as e:
                            logger.debug(f"Could not parse app properties: {e}")
                            
        except Exception as e:
            logger.debug(f"Could not extract Office metadata: {e}")
        
        return metadata


class WordHandler(BaseOfficeHandler):
    """Handler for Word documents (.docx, .doc)"""
    
    def can_handle(self, file_extension: str) -> bool:
        """Check if this is a Word document"""
        return file_extension.lower() in ['.docx', '.doc']
    
    def process(self, file_path: Path, metadata: Dict[str, Any], mode: str = 'intro') -> Dict[str, Any]:
        """
        Process Word document
        
        Args:
            file_path: Path to Word file
            metadata: Basic file metadata
            mode: 'intro' (first 10k chars) or 'full' (all content)
            
        Returns:
            Dictionary with content, metadata, and error
        """
        # Handle .docx files
        if file_path.suffix.lower() == '.docx':
            if not DOCX_AVAILABLE:
                return {
                    'content': '',
                    'metadata': metadata,
                    'error': 'python-docx not installed. Install with: pip install python-docx'
                }
            
            try:
                return self._process_docx(file_path, metadata, mode)
            except Exception as e:
                logger.error(f"Error processing DOCX: {e}")
                return {
                    'content': '',
                    'metadata': metadata,
                    'error': f'Failed to process DOCX: {str(e)}'
                }
        
        # Handle .doc files
        elif file_path.suffix.lower() == '.doc':
            if DOC2TXT_AVAILABLE:
                try:
                    content = doc2txt.process(str(file_path))
                    if mode == 'intro' and content:
                        content = content[:self.intro_chars]
                    
                    metadata['extraction_method'] = 'doc2txt'
                    return {
                        'content': content or '',
                        'metadata': metadata,
                        'error': None
                    }
                except Exception as e:
                    logger.error(f"Error with doc2txt: {e}")
            
            # Fallback for .doc files
            return {
                'content': '',
                'metadata': metadata,
                'error': 'Legacy .doc format. Convert to .docx or install doc2txt'
            }
    
    def _process_docx(self, file_path: Path, metadata: Dict[str, Any], mode: str) -> Dict[str, Any]:
        """Process modern .docx files"""
        doc = Document(str(file_path))
        
        # Extract text from paragraphs
        paragraphs = []
        char_count = 0
        
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
                char_count += len(para.text)
                
                # For intro mode, stop at character limit
                if mode == 'intro' and char_count >= self.intro_chars:
                    break
        
        # Extract text from tables if needed
        if mode == 'full' or char_count < self.intro_chars:
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    
                    if row_text:
                        table_row = ' | '.join(row_text)
                        paragraphs.append(table_row)
                        char_count += len(table_row)
                        
                        if mode == 'intro' and char_count >= self.intro_chars:
                            break
                
                if mode == 'intro' and char_count >= self.intro_chars:
                    break
        
        content = '\n\n'.join(paragraphs)
        
        # Update metadata
        metadata['paragraphs'] = len(doc.paragraphs)
        metadata['tables'] = len(doc.tables)
        metadata['extraction_method'] = 'python-docx'
        
        # Add Office metadata
        office_metadata = self._extract_basic_metadata(file_path)
        metadata.update(office_metadata)
        
        return {
            'content': content,
            'metadata': metadata,
            'error': None
        }


class ExcelHandler(BaseOfficeHandler):
    """Handler for Excel documents (.xlsx, .xls)"""
    
    def __init__(self):
        """Initialize handler"""
        super().__init__()
        self.intro_sheets = 5  # Number of sheets for intro mode
        self.max_rows_intro = 100  # Max rows per sheet in intro mode
    
    def can_handle(self, file_extension: str) -> bool:
        """Check if this is an Excel document"""
        return file_extension.lower() in ['.xlsx', '.xls']
    
    def process(self, file_path: Path, metadata: Dict[str, Any], mode: str = 'intro') -> Dict[str, Any]:
        """
        Process Excel document
        
        Args:
            file_path: Path to Excel file
            metadata: Basic file metadata
            mode: 'intro' (first 5 sheets, 100 rows each) or 'full' (all content)
            
        Returns:
            Dictionary with content, metadata, and error
        """
        if not OPENPYXL_AVAILABLE:
            return {
                'content': '',
                'metadata': metadata,
                'error': 'openpyxl not installed. Install with: pip install openpyxl'
            }
        
        try:
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            
            content_parts = []
            sheet_count = 0
            total_rows = 0
            
            # Determine sheets to process
            sheets_to_process = wb.sheetnames
            if mode == 'intro':
                sheets_to_process = sheets_to_process[:self.intro_sheets]
            
            for sheet_name in sheets_to_process:
                sheet = wb[sheet_name]
                sheet_content = [f"=== Sheet: {sheet_name} ==="]
                sheet_rows = 0
                
                # Determine row limit
                max_rows = None if mode == 'full' else self.max_rows_intro
                
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                    if max_rows and row_idx >= max_rows:
                        break
                    
                    # Filter out empty rows
                    if any(cell is not None for cell in row):
                        row_values = [str(cell) if cell is not None else '' for cell in row]
                        sheet_content.append('\t'.join(row_values))
                        sheet_rows += 1
                
                if sheet_rows > 0:
                    content_parts.append('\n'.join(sheet_content))
                    total_rows += sheet_rows
                
                sheet_count += 1
            
            wb.close()
            
            content = '\n\n'.join(content_parts)
            
            # Update metadata
            metadata['sheets'] = len(wb.sheetnames)
            metadata['sheets_processed'] = sheet_count
            metadata['rows_extracted'] = total_rows
            metadata['extraction_method'] = 'openpyxl'
            
            # Add Office metadata
            office_metadata = self._extract_basic_metadata(file_path)
            metadata.update(office_metadata)
            
            return {
                'content': content,
                'metadata': metadata,
                'error': None
            }
            
        except Exception as e:
            logger.error(f"Error processing Excel: {e}")
            return {
                'content': '',
                'metadata': metadata,
                'error': f'Failed to process Excel: {str(e)}'
            }


class PowerPointHandler(BaseOfficeHandler):
    """Handler for PowerPoint documents (.pptx, .ppt)"""
    
    def __init__(self):
        """Initialize handler"""
        super().__init__()
        self.intro_slides = 10  # Number of slides for intro mode
    
    def can_handle(self, file_extension: str) -> bool:
        """Check if this is a PowerPoint document"""
        return file_extension.lower() in ['.pptx', '.ppt']
    
    def process(self, file_path: Path, metadata: Dict[str, Any], mode: str = 'intro') -> Dict[str, Any]:
        """
        Process PowerPoint document
        
        Args:
            file_path: Path to PowerPoint file
            metadata: Basic file metadata
            mode: 'intro' (first 10 slides) or 'full' (all slides)
            
        Returns:
            Dictionary with content, metadata, and error
        """
        if not PPTX_AVAILABLE:
            return {
                'content': '',
                'metadata': metadata,
                'error': 'python-pptx not installed. Install with: pip install python-pptx'
            }
        
        # Only handle .pptx files
        if file_path.suffix.lower() != '.pptx':
            return {
                'content': '',
                'metadata': metadata,
                'error': 'Legacy .ppt format not supported. Please convert to .pptx'
            }
        
        try:
            prs = Presentation(str(file_path))
            
            content_parts = []
            slides_processed = 0
            
            # Determine slides to process
            max_slides = len(prs.slides) if mode == 'full' else min(self.intro_slides, len(prs.slides))
            
            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx >= max_slides:
                    break
                
                slide_content = [f"=== Slide {slide_idx + 1} ==="]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                    
                    # Extract text from tables
                    if shape.has_table:
                        table_text = []
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                table_text.append(' | '.join(row_text))
                        if table_text:
                            slide_content.extend(table_text)
                
                # Extract notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    slide_content.append(f"[Notes: {slide.notes_slide.notes_text_frame.text.strip()}]")
                
                if len(slide_content) > 1:  # More than just the slide header
                    content_parts.append('\n'.join(slide_content))
                    slides_processed += 1
            
            content = '\n\n'.join(content_parts)
            
            # Update metadata
            metadata['total_slides'] = len(prs.slides)
            metadata['slides_processed'] = slides_processed
            metadata['extraction_method'] = 'python-pptx'
            
            # Add Office metadata
            office_metadata = self._extract_basic_metadata(file_path)
            metadata.update(office_metadata)
            
            return {
                'content': content,
                'metadata': metadata,
                'error': None
            }
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint: {e}")
            return {
                'content': '',
                'metadata': metadata,
                'error': f'Failed to process PowerPoint: {str(e)}'
            }


# Compatibility exports
__all__ = ['WordHandler', 'ExcelHandler', 'PowerPointHandler']